from pptx import Presentation
import os
import xlrd
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR_INDEX

from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import pandas as pd

df_xlsx = pd.read_excel("D:/Upthrust/Frank/PoC Cookbook Automation/Excel/Volkswagen Cookbook Excel.xlsx")


prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]


title.text = "VW Cookbook"
subtitle.text = "Growthmarketing"

# Voor uitleg van de x, y, i kijk onderaan
def main(x, y, i):

    z = df_xlsx.iloc[x, y]
    normal_slide = prs.slide_layouts[5]
    slide_2 = prs.slides.add_slide(normal_slide)
    title = slide_2.shapes.title

    title.text = "Experiment 1: Search Campaign Per Model"

    for shape in slide_2.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        # do things with the text frame


    text_frame = shape.text_frame
    text_frame.clear()

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(25.3)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme



    txBox = slide_2.shapes.add_textbox(Inches(0), Inches(6.5), Inches(2), Inches(1))
    tf = txBox.text_frame



    for shape in slide_2.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        # do things with the text frame


    text_frame = shape.text_frame

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Started: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme

    y = y + 1
    z = str(df_xlsx.iloc[x,y])


    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme

    p = tf.add_paragraph()



    run = p.add_run()
    run.text = 'Status: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme

    y = y + 1
    z = df_xlsx.iloc[x,y]

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme


    p = tf.add_paragraph()

    run = p.add_run()
    run.text = 'Channel(s): '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme


    y = y + 1
    z = df_xlsx.iloc[x,y]

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme


    txBox = slide_2.shapes.add_textbox(Inches(3), Inches(6.5), Inches(2), Inches(1))
    tf = txBox.text_frame

    for shape in slide_2.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        # do things with the text frame


    text_frame = shape.text_frame

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Current Results: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme

    txBox = slide_2.shapes.add_textbox(Inches(5), Inches(6.5), Inches(2), Inches(1))
    tf = txBox.text_frame


    for shape in slide_2.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        # do things with the text frame


    text_frame = shape.text_frame

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Total Reach: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme


    text_frame = shape.text_frame

    y = y + 1
    z = str(df_xlsx.iloc[x,y])

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme

    p = tf.add_paragraph()

    text_frame = shape.text_frame


    run = p.add_run()
    run.text = 'Total Clicks: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme


    text_frame = shape.text_frame

    y = y + 1
    z = str(df_xlsx.iloc[x,y])


    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme

    p = tf.add_paragraph()

    text_frame = shape.text_frame


    run = p.add_run()
    run.text = 'Media Spend: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme


    text_frame = shape.text_frame

    y = y + 1
    z = str(df_xlsx.iloc[x,y])

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme


    txBox = slide_2.shapes.add_textbox(Inches(7.5), Inches(6.5), Inches(2), Inches(1))
    tf = txBox.text_frame


    for shape in slide_2.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        # do things with the text frame


    text_frame = shape.text_frame

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Total Leads: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme


    text_frame = shape.text_frame

    y = y + 1
    z = str(df_xlsx.iloc[x,y])

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme

    p = tf.add_paragraph()

    text_frame = shape.text_frame


    run = p.add_run()
    run.text = 'Cost Per Lead: : '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme


    text_frame = shape.text_frame

    y = y + 1
    z = str(df_xlsx.iloc[x,y])

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme


    p = tf.add_paragraph()

    text_frame = shape.text_frame


    run = p.add_run()
    run.text = 'Ads Set Up: : '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme

    y = y + 1
    z = str(df_xlsx.iloc[x,y])

    text_frame = shape.text_frame

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme
    prs.save('D:/Upthrust/Frank/PoC Cookbook Automation/Powerpoints/Test.pptx')

    x = x + 1
    y = 0
    z = str(df_xlsx.iloc[x,y])
    print(x, y, z)

    # Dit is wat de loop (eigenlijk een recursieve functie, een functie die zichzelf aan blijft roepen) draaiend houdt. Je had dit eerder ook al, maar toen had je
    # nog geen "limiet" ingebouwd waardoor het programma draaide totdat er een error was. Nu stopt het programma netjes zodra de variabele i 0 is
    i -= 1
    if i > 0:
      main(x, y, i)


# the row
x = 0

# the column
y = 0

# Als je python code schrijft, en je maakt een main() functie, wordt deze niet automatisch gedraait. De code hieronder is een code waarbij de main() function wordt gedraait als
# de python code direct wordt gerunned. Reference: https://stackoverflow.com/questions/419163/what-does-if-name-main-do
if __name__ == "__main__":
  # Je moet x en y meegeven, aangezien dit local variables zijn. de i (hier 25) die meegegeven is hoevaak het programma draait.
  main(x, y, 25)

os.startfile("D:/Upthrust/Frank/PoC Cookbook Automation/Powerpoints/Test.pptx")


# Wat je nog zou kunnen doen is kritisch naar de code in main() kijken. Gebeuren er soms dingen die vaker herhaald worden, kunnen deze misschien in een functie
# worden gezet, waardoor je minder code zou hebben? Een van de belangerijkste regels in code is dat als je hetzelfde vaker doet, je er waarschijnlijk een
# aanroepbare functie/class/whatever van kan maken. Hierdoor krijg je ook geen bestanden van 400 regels die lastig te lezen zijn.